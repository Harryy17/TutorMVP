"""
Document Processor — Multi-stage Document Processing Pipeline with Fast-Path & Background Enrichment.

Stages:
[Stage 0] Save file, register doc_id, status = "processing_text"
[Stage 1 — FAST PATH, ~seconds]
  - Extract text per page
  - Chunk into 500-800 characters with ~15% overlap
  - Store chunks with metadata: { doc_id, page, source_type: "text" }
  - status = "text_ready" (User can ask questions immediately)
[Stage 2 — TABLE EXTRACTION, background task]
  - pdfplumber.extract_tables()
  - Convert tables to Markdown strings
  - Store chunks with metadata: { doc_id, page, source_type: "table" }
[Stage 3 — IMAGE EXTRACTION & CAPTIONING, background task]
  - Extract embedded images from PDF pages
  - Vision-caption each image using Gemini VLM
  - Store captions with metadata: { doc_id, page, source_type: "image_caption" }
  - status = "fully_processed"
"""
import os
import io
import re
import asyncio
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
from dataclasses import dataclass, field

import pypdf
import pdfplumber
from PIL import Image

from app.rag.vlm_client import GeminiVLMClient, _get_active_gemini_key
from app.rag.sqlite_fts_store import sqlite_fts_store, get_session_store
import httpx



@dataclass
class DocumentChunk:
    chunk_id: str
    doc_id: str
    page: int
    source_type: str  # "text" | "table" | "image_caption"
    content: str
    metadata: Dict[str, Any] = field(default_factory=dict)


@dataclass
class DocumentRecord:
    doc_id: str
    file_path: str
    file_name: str
    subject: str
    session_id: Optional[str] = None
    status: str = "processing_text"  # "processing_text" | "text_ready" | "processing_enrichment" | "fully_processed" | "error"
    chunks: List[DocumentChunk] = field(default_factory=list)
    stats: Dict[str, int] = field(default_factory=lambda: {"text_chunks": 0, "tables": 0, "images": 0})
    error_message: Optional[str] = None



class DocumentProcessor:
    """
    Manages document ingestion, chunking, background table/image extraction,
    and metadata-tagged context retrieval.
    """

    def __init__(self):
        self._docs: Dict[str, DocumentRecord] = {}
        self.vlm = GeminiVLMClient()

    def get_document(self, doc_id: str) -> Optional[DocumentRecord]:
        return self._docs.get(doc_id)

    # ─── STAGE 1: INGESTION (Fast Text + VLM OCR for Scanned/Images) ─────
    async def ingest_document(
        self,
        doc_id: str,
        file_path: str,
        file_name: str,
        subject: str = "General",
        session_id: Optional[str] = None,
    ) -> DocumentRecord:
        """
        Stage 1: Multi-modal document ingestion.
        - Plain text/code: Direct read & chunk.
        - Digital PDF: Fast text extraction via pypdf.
        - Scanned PDF / Image: OCR & transcription via Google Gemini VLM.
        """
        doc = DocumentRecord(
            doc_id=doc_id,
            file_path=file_path,
            file_name=file_name,
            subject=subject,
            session_id=session_id,
            status="processing_text",
        )
        self._docs[doc_id] = doc

        ext = Path(file_path).suffix.lower()

        # 1. Plain text / code / markdown files
        if ext in {".txt", ".md", ".py", ".csv", ".json"}:
            try:
                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    full_text = f.read()
                chunks = self._chunk_text(full_text, doc_id=doc_id, page=1, source_type="text")
                doc.chunks.extend(chunks)
                doc.stats["text_chunks"] = len(chunks)
                doc.status = "fully_processed"
                store = get_session_store(session_id)
                store.index_chunks([
                    {"chunk_id": c.chunk_id, "doc_id": c.doc_id, "page": c.page, "source_type": c.source_type, "content": c.content}
                    for c in chunks
                ])
                return doc
            except Exception as e:
                doc.status = "error"
                doc.error_message = str(e)
                return doc

        # 2. Standalone image files (.png, .jpg, .jpeg, .webp, .bmp, .tiff)
        if ext in {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tiff", ".tif"}:
            try:
                print(f"[DocProcessor] Image file detected: {file_name}. Transcribing using Gemini VLM...")
                with open(file_path, "rb") as f:
                    img_bytes = f.read()

                mime = "image/png" if ext == ".png" else "image/jpeg"
                vlm_text = await self.vlm.extract_text_from_image(
                    img_bytes, mime_type=mime, context_hint=subject
                )

                if vlm_text and vlm_text.strip():
                    chunks = self._chunk_text(vlm_text, doc_id=doc_id, page=1, source_type="text")
                    doc.chunks.extend(chunks)
                    doc.stats["text_chunks"] = len(chunks)
                    store = get_session_store(session_id)
                    store.index_chunks([
                        {"chunk_id": c.chunk_id, "doc_id": c.doc_id, "page": c.page, "source_type": c.source_type, "content": c.content}
                        for c in chunks
                    ])
                    print(f"[DocProcessor] Successfully extracted {len(chunks)} text chunks from image with VLM.")
                else:
                    print(f"[DocProcessor] VLM returned empty text for image {file_name}")

                doc.status = "fully_processed"
                return doc
            except Exception as e:
                print(f"[DocProcessor] Image VLM extraction error: {e}")
                doc.status = "error"
                doc.error_message = str(e)
                return doc

        # 3. PDF documents (digital or scanned)
        if ext == ".pdf":
            try:
                reader = pypdf.PdfReader(file_path)
                total_pages = len(reader.pages)
                text_chunks: List[DocumentChunk] = []
                scanned_pages_to_vlm: List[int] = []

                # Digital text pass: extract and index all digital text pages immediately (<1.5s total)
                for page_idx in range(total_pages):
                    page_num = page_idx + 1
                    try:
                        page_text = reader.pages[page_idx].extract_text() or ""
                    except Exception:
                        page_text = ""

                    # If page text is rich enough, chunk it directly
                    if len(page_text.strip()) >= 40:
                        chunks = self._chunk_text(
                            page_text,
                            doc_id=doc_id,
                            page=page_num,
                            source_type="text"
                        )
                        text_chunks.extend(chunks)
                    else:
                        # Page has no or negligible text -> Scanned page!
                        scanned_pages_to_vlm.append(page_idx)



                # If scanned pages exist, use Gemini VLM to extract text from images in PARALLEL
                if scanned_pages_to_vlm:
                    fast_scanned = scanned_pages_to_vlm[:6]
                    print(f"[DocProcessor] Scanned PDF detected ({len(scanned_pages_to_vlm)}/{total_pages} scanned pages). Running PARALLEL Gemini VLM OCR on {len(fast_scanned)} pages...")

                    async def _ocr_single_page(p_idx: int) -> List[DocumentChunk]:
                        p_num = p_idx + 1
                        try:
                            p_img = await asyncio.to_thread(self.vlm.render_pdf_page_to_image, file_path, p_idx, 120)
                            if not p_img:
                                return []
                            p_ocr_text = await self.vlm.extract_text_from_image(
                                p_img, mime_type="image/png", context_hint=subject
                            )
                            if p_ocr_text and p_ocr_text.strip():
                                p_chunks = self._chunk_text(
                                    p_ocr_text,
                                    doc_id=doc_id,
                                    page=p_num,
                                    source_type="text"
                                )
                                print(f"[DocProcessor] VLM transcribed scanned PDF Page {p_num} -> {len(p_chunks)} chunks.")
                                return p_chunks
                        except Exception as ocr_err:
                            print(f"[DocProcessor] Error on parallel page {p_num}: {ocr_err}")
                        return []

                    # Execute all pages concurrently in parallel
                    page_results = await asyncio.gather(*[_ocr_single_page(p) for p in fast_scanned], return_exceptions=True)
                    for res in page_results:
                        if isinstance(res, list):
                            text_chunks.extend(res)


                doc.chunks.extend(text_chunks)
                doc.stats["text_chunks"] = len(text_chunks)
                doc.status = "text_ready"

                # Persist to local session-isolated SQLite FTS5 index
                store = get_session_store(session_id)
                store.index_chunks([
                    {"chunk_id": c.chunk_id, "doc_id": c.doc_id, "page": c.page, "source_type": c.source_type, "content": c.content}
                    for c in text_chunks
                ])
                print(f"[DocProcessor] Fast path completed: {len(text_chunks)} text chunks indexed for doc {doc_id}.")
            except Exception as e:
                print(f"[DocProcessor] Fast path text extraction error: {e}")
                doc.status = "text_ready"

        # 4. Word documents (.docx, .doc)
        elif ext in {".docx", ".doc"}:
            try:
                import docx
                doc_file = docx.Document(file_path)
                full_paragraphs = []
                for p in doc_file.paragraphs:
                    txt = p.text.strip()
                    if txt:
                        full_paragraphs.append(txt)
                
                # Also extract tables inside docx
                for tbl in doc_file.tables:
                    for row in tbl.rows:
                        row_text = " | ".join([cell.text.strip() for cell in row.cells if cell.text.strip()])
                        if row_text:
                            full_paragraphs.append(row_text)

                full_text = "\n\n".join(full_paragraphs)
                if not full_text.strip():
                    raise ValueError(f"No readable text extracted from {file_name}")

                chunks = self._chunk_text(full_text, doc_id=doc_id, page=1, source_type="text")
                doc.chunks.extend(chunks)
                doc.stats["text_chunks"] = len(chunks)
                doc.status = "fully_processed"
                store = get_session_store(session_id)
                store.index_chunks([
                    {"chunk_id": c.chunk_id, "doc_id": c.doc_id, "page": c.page, "source_type": c.source_type, "content": c.content}
                    for c in chunks
                ])
                print(f"[DocProcessor] Word doc ingestion completed: {len(chunks)} text chunks indexed for {file_name}.")
                return doc
            except Exception as e:
                print(f"[DocProcessor] DOCX extraction error: {e}")
                doc.status = "error"
                doc.error_message = str(e)
                return doc

        # 5. PowerPoint presentations (.pptx, .ppt)
        elif ext in {".pptx", ".ppt"}:
            try:
                from pptx import Presentation
                prs = Presentation(file_path)
                slide_chunks = []
                for slide_idx, slide in enumerate(prs.slides):
                    slide_texts = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_texts.append(shape.text.strip())
                    if slide_texts:
                        slide_content = "\n".join(slide_texts)
                        chunks = self._chunk_text(slide_content, doc_id=doc_id, page=slide_idx + 1, source_type="text")
                        slide_chunks.extend(chunks)

                if not slide_chunks:
                    raise ValueError(f"No readable text found in presentation {file_name}")

                doc.chunks.extend(slide_chunks)
                doc.stats["text_chunks"] = len(slide_chunks)
                doc.status = "fully_processed"
                store = get_session_store(session_id)
                store.index_chunks([
                    {"chunk_id": c.chunk_id, "doc_id": c.doc_id, "page": c.page, "source_type": c.source_type, "content": c.content}
                    for c in slide_chunks
                ])
                print(f"[DocProcessor] PPTX ingestion completed: {len(slide_chunks)} chunks indexed for {file_name}.")
                return doc
            except Exception as e:
                print(f"[DocProcessor] PPTX extraction error: {e}")
                doc.status = "error"
                doc.error_message = str(e)
                return doc

        return doc

    def get_document_text(self, doc_id: str, max_chars: int = 12000) -> str:
        """
        Retrieves clean assembled text of the document across all chunks
        for feeding into the LLM/Teaching Engine/Exam Engine.
        """
        doc = self._docs.get(doc_id)
        if doc and doc.chunks:
            text_pieces = [c.content for c in doc.chunks if c.source_type == "text"]
            if not text_pieces:
                text_pieces = [c.content for c in doc.chunks]
            full = "\n\n".join(text_pieces)
            return full[:max_chars]

        # Fallback to session store SQLite search
        store = get_session_store(doc_id)
        matches = store.search(doc_id=doc_id, query="*", limit=20)
        if matches:
            return "\n\n".join(m["content"] for m in matches)[:max_chars]
        return ""

    # ─── STAGE 2 & 3: ASYNC BACKGROUND ENRICHMENT ─────────────────────────
    async def run_background_enrichment(self, doc_id: str):
        """
        Non-blocking background task to:
        1. OCR any remaining scanned PDF pages beyond page 5
        2. Extract tables (Stage 2)
        3. Extract & caption diagram images with Gemini VLM (Stage 3)
        Appends chunks directly to the active document and local SQLite FTS5 store.
        """
        doc = self._docs.get(doc_id)
        if not doc or not os.path.exists(doc.file_path):
            return

        ext = Path(doc.file_path).suffix.lower()
        if ext != ".pdf":
            doc.status = "fully_processed"
            return

        doc.status = "processing_enrichment"



        # ── Stage 2: Table Extraction ──
        try:
            table_chunks = await asyncio.to_thread(self._extract_tables_from_pdf, doc.file_path, doc_id)
            if table_chunks:
                doc.chunks.extend(table_chunks)
                doc.stats["tables"] = len(table_chunks)
                store = get_session_store(doc.session_id)
                store.index_chunks([
                    {"chunk_id": c.chunk_id, "doc_id": c.doc_id, "page": c.page, "source_type": c.source_type, "content": c.content}
                    for c in table_chunks
                ])
                print(f"[DocProcessor] Stage 2 Complete: Extracted & indexed {len(table_chunks)} tables for doc {doc_id}")
        except Exception as e:
            print(f"[DocProcessor] Stage 2 Table Extraction warning: {e}")

        # ── Stage 3: Image Extraction & VLM Captioning ──
        try:
            image_chunks = await self._extract_and_caption_images(doc.file_path, doc_id)
            if image_chunks:
                doc.chunks.extend(image_chunks)
                doc.stats["images"] = len(image_chunks)
                store = get_session_store(doc.session_id)
                store.index_chunks([
                    {"chunk_id": c.chunk_id, "doc_id": c.doc_id, "page": c.page, "source_type": c.source_type, "content": c.content}
                    for c in image_chunks
                ])
                print(f"[DocProcessor] Stage 3 Complete: Extracted & captioned {len(image_chunks)} images for doc {doc_id}")
        except Exception as e:
            print(f"[DocProcessor] Stage 3 Image Captioning warning: {e}")

        doc.status = "fully_processed"
        print(f"[DocProcessor] Document {doc_id} is now FULLY_PROCESSED. Total chunks in local SQLite FTS5: {len(doc.chunks)} (Text: {doc.stats['text_chunks']}, Tables: {doc.stats['tables']}, Images: {doc.stats['images']})")


    # ─── TABLE EXTRACTION HELPER ──────────────────────────────────────────
    def _extract_tables_from_pdf(self, file_path: str, doc_id: str) -> List[DocumentChunk]:
        """Extracts tables per page using pdfplumber and formats them as Markdown tables."""
        table_chunks: List[DocumentChunk] = []
        try:
            with pdfplumber.open(file_path) as pdf:
                for page_idx, page in enumerate(pdf.pages):
                    page_num = page_idx + 1
                    tables = page.extract_tables()
                    for t_idx, table in enumerate(tables):
                        if not table or len(table) < 2:
                            continue

                        # Clean and format into markdown table
                        md_table = self._format_table_as_markdown(table)
                        if md_table.strip():
                            chunk = DocumentChunk(
                                chunk_id=f"{doc_id}_p{page_num}_tbl_{t_idx + 1}",
                                doc_id=doc_id,
                                page=page_num,
                                source_type="table",
                                content=md_table,
                                metadata={"table_index": t_idx + 1, "rows": len(table), "cols": len(table[0]) if table else 0}
                            )
                            table_chunks.append(chunk)
        except Exception as e:
            print(f"[DocProcessor] pdfplumber table error: {e}")
        return table_chunks

    def _format_table_as_markdown(self, table: List[List[Any]]) -> str:
        """Converts raw table rows into clean, structured Markdown table."""
        if not table:
            return ""

        # Normalize cells
        cleaned_rows = []
        for row in table:
            cleaned_row = [str(cell).replace("\n", " ").strip() if cell is not None else "" for cell in row]
            cleaned_rows.append(cleaned_row)

        if not cleaned_rows:
            return ""

        headers = cleaned_rows[0]
        # Ensure header is not all empty
        if not any(headers):
            headers = [f"Col {i+1}" for i in range(len(headers))]

        col_count = len(headers)
        md_lines = []
        md_lines.append("| " + " | ".join(headers) + " |")
        md_lines.append("| " + " | ".join(["---"] * col_count) + " |")

        for row in cleaned_rows[1:]:
            # Pad or truncate row to match col_count
            padded = row + [""] * (col_count - len(row)) if len(row) < col_count else row[:col_count]
            md_lines.append("| " + " | ".join(padded) + " |")

        return "\n".join(md_lines)

    # ─── IMAGE EXTRACTION & VLM CAPTIONING HELPER ─────────────────────────
    async def _extract_and_caption_images(self, file_path: str, doc_id: str, max_images: int = 8) -> List[DocumentChunk]:
        """Extracts images from PDF pages and runs Gemini VLM factual captioning with PyMuPDF."""
        image_chunks: List[DocumentChunk] = []

        try:
            try:
                import fitz  # PyMuPDF
                pdf_doc = fitz.open(file_path)
                for page_num in range(len(pdf_doc)):
                    if len(image_chunks) >= max_images:
                        break
                    page = pdf_doc[page_num]
                    page_text = page.get_text()
                    image_list = page.get_images(full=True)

                    for img_index, img in enumerate(image_list):
                        if len(image_chunks) >= max_images:
                            break
                        xref = img[0]
                        base_image = pdf_doc.extract_image(xref)
                        image_bytes = base_image.get("image")
                        image_ext = base_image.get("ext", "jpeg")

                        if not image_bytes or len(image_bytes) < 3000:
                            continue

                        # Extract potential figure numbers and captions from page text (e.g. "Fig. 4.3: Bamboo plantation in North-east")
                        fig_matches = re.findall(r"(?:Fig(?:ure)?\.?\s*\d+(?:\.\d+)?[:\s\-\–][^\n\.\;]{4,100})", page_text, re.IGNORECASE)
                        page_fig_context = "; ".join(fig_matches[:3]) if fig_matches else ""

                        caption = await self._caption_image_with_vlm(image_bytes, mime_type=f"image/{image_ext}", context_hint=page_fig_context)
                        await asyncio.sleep(0.2)  # Yield loop to allow user chat requests priority

                        content_parts = [f"Figure / Diagram on Page {page_num + 1}"]
                        if page_fig_context:
                            content_parts.append(f"Figure Caption in text: {page_fig_context}")
                        if caption:
                            content_parts.append(f"Visual Details: {caption.strip()}")

                        if len(content_parts) > 1:
                            chunk = DocumentChunk(
                                chunk_id=f"{doc_id}_p{page_num+1}_img_{img_index+1}",
                                doc_id=doc_id,
                                page=page_num + 1,
                                source_type="image_caption",
                                content=". ".join(content_parts),
                                metadata={"image_index": img_index + 1}
                            )
                            image_chunks.append(chunk)
            except Exception as fitz_err:
                print(f"[DocProcessor] PyMuPDF extraction fallback to pypdf: {fitz_err}")
                reader = pypdf.PdfReader(file_path)
                for page_idx, page in enumerate(reader.pages):
                    page_num = page_idx + 1
                    if len(image_chunks) >= max_images:
                        break
                    try:
                        page_images = list(page.images)
                    except Exception:
                        page_images = []
                    for img_idx, img_obj in enumerate(page_images):
                        if len(image_chunks) >= max_images:
                            break
                        img_bytes = img_obj.data
                        if len(img_bytes) < 3000:
                            continue
                        caption = await self._caption_image_with_vlm(img_bytes)
                        if caption and len(caption.strip()) > 10:
                            chunk = DocumentChunk(
                                chunk_id=f"{doc_id}_p{page_num}_img_{img_idx + 1}",
                                doc_id=doc_id,
                                page=page_num,
                                source_type="image_caption",
                                content=f"Figure/Diagram on Page {page_num}: {caption.strip()}",
                            )
                            image_chunks.append(chunk)
        except Exception as e:
            print(f"[DocProcessor] PDF Image extraction error: {e}")

        return image_chunks

    async def _caption_image_with_vlm(self, image_bytes: bytes, mime_type: str = "image/jpeg", context_hint: str = "") -> str:
        """Calls Gemini Vision API with cascade to describe figures and diagrams factually."""
        api_key = _get_active_gemini_key()
        if not api_key:
            return ""

        import base64
        b64_data = base64.b64encode(image_bytes).decode("utf-8")

        prompt = (
            "Analyze and describe this educational figure, photo, diagram, or chart in detail.\n"
            f"{f'Context / Surrounding Caption: {context_hint}' if context_hint else ''}\n"
            "MANDATORY:\n"
            "1. State the figure title, number (e.g. Fig. 1.2), or subject matter.\n"
            "2. Describe the key visual elements, features, processes, and any text labels or data shown.\n"
            "3. Keep the description clear, factual, and concise."
        )

        payload = {
            "contents": [
                {
                    "parts": [
                        {"text": prompt},
                        {
                            "inline_data": {
                                "mime_type": mime_type if "image/" in mime_type else "image/jpeg",
                                "data": b64_data,
                            }
                        }
                    ]
                }
            ],
            "generationConfig": {"temperature": 0.1, "maxOutputTokens": 400}
        }

        models_to_try = [
            "gemini-3.5-flash-lite",
            "gemini-3.6-flash",
            "gemini-flash-latest",
            "gemini-3.7-flash",
            "gemini-3.5-flash",
        ]
        for model in models_to_try:
            url = f"https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent?key={api_key}"
            try:
                async with httpx.AsyncClient(timeout=15.0) as client:
                    res = await client.post(url, json=payload)
                    if res.status_code == 200:
                        data = res.json()
                        candidates = data.get("candidates", [])
                        if candidates:
                            parts = candidates[0].get("content", {}).get("parts", [])
                            if parts:
                                return parts[0].get("text", "").strip()
            except Exception as e:
                continue

        return ""

    # ─── CHUNKING UTILITY ─────────────────────────────────────────────────
    def _chunk_text(
        self,
        text: str,
        doc_id: str,
        page: int,
        source_type: str = "text",
        chunk_size: int = 700,
        overlap: int = 100,
    ) -> List[DocumentChunk]:
        """Chunks text with 500-800 character windows and ~15% overlap."""
        cleaned = re.sub(r"\s+", " ", text).strip()
        if not cleaned:
            return []

        chunks: List[DocumentChunk] = []
        start = 0
        chunk_idx = 1

        while start < len(cleaned):
            end = start + chunk_size
            chunk_content = cleaned[start:end].strip()

            if chunk_content:
                chunk = DocumentChunk(
                    chunk_id=f"{doc_id}_p{page}_c{chunk_idx}",
                    doc_id=doc_id,
                    page=page,
                    source_type=source_type,
                    content=chunk_content,
                )
                chunks.append(chunk)
                chunk_idx += 1

            if end >= len(cleaned):
                break
            start += chunk_size - overlap

        return chunks

    # ─── QUERY RETRIEVAL FLOW ─────────────────────────────────────────────
    def retrieve_context(
        self,
        doc_id: str,
        query: str,
        top_k: int = 6,
        session_id: Optional[str] = None,
    ) -> Tuple[str, str, Dict[str, Any]]:

        """
        Retrieves top-k relevant chunks across ALL source_types (text, table, image_caption),
        with specialized routing for table-related queries.

        Returns:
          - formatted_context_block (with source_type and page tags)
          - doc_status_note (if still processing background enrichment)
          - debug_metadata
        """
        doc = self._docs.get(doc_id)
        store = get_session_store(session_id or doc_id or (getattr(doc, "session_id", None) if doc else None))

        query_lower = query.lower().strip()
        table_keywords = {"table", "value", "compare", "how many", "number", "data", "columns", "rows", "statistic", "percent", "metric", "versus", "vs"}
        is_table_query = any(re.search(rf"\b{re.escape(kw)}\b", query_lower) for kw in table_keywords)

        # Detect specific Page Number e.g. "page 22", "page number 22", "pg 22", "p. 22"
        page_match = re.search(r"\b(?:page|pg|p\.?)\s*(?:number|no\.?)?\s*(\d+)\b", query_lower)
        target_page = int(page_match.group(1)) if page_match else None

        # Detect specific Figure / Table identifiers e.g. "fig. 3.4", "figure 3.4", "fig 2", "table 1.1"
        fig_refs = re.findall(r"\b(?:fig(?:ure)?\.?\s*\d+(?:\.\d+)?|table\s*\d+(?:\.\d+)?)\b", query_lower)
        clean_fig_patterns = []
        for ref in fig_refs:
            num_match = re.search(r"\d+(?:\.\d+)?", ref)
            if num_match:
                num = num_match.group(0)
                clean_fig_patterns.append(rf"\bfig(?:ure)?\.?\s*{re.escape(num)}\b")

        # Tokenize query including numbers and short domain terms
        query_words = set(re.findall(r"[a-z0-9]+(?:\.[a-z0-9]+)?", query_lower))
        selected_chunks: List[DocumentChunk] = []

        # 1. If target page is specified, fetch all chunks from that exact page directly
        if target_page is not None:
            page_rows = store.get_page_chunks(doc_id=doc_id, page=target_page)
            for r in page_rows:
                selected_chunks.append(
                    DocumentChunk(
                        chunk_id=r["chunk_id"],
                        doc_id=r["doc_id"],
                        page=r["page"],
                        source_type=r["source_type"],
                        content=r["content"],
                    )
                )

        # 2. In-memory scoring pass if chunks are available in RAM
        if doc and doc.chunks:
            scored_chunks: List[Tuple[float, DocumentChunk]] = []
            for chunk in doc.chunks:
                score = 0.0
                content_lower = chunk.content.lower()

                # Exact page boost
                if target_page is not None and chunk.page == target_page:
                    score += 60.0
                elif target_page is not None and abs(chunk.page - target_page) == 1:
                    score += 15.0

                # Exact Figure / Table reference match
                for pat in clean_fig_patterns:
                    if re.search(pat, content_lower):
                        score += 35.0

                # Keyword matching score
                for word in query_words:
                    count = content_lower.count(word)
                    if count > 0:
                        score += 1.0 + min(count * 0.5, 3.0)

                # Boost table chunks
                if is_table_query and chunk.source_type == "table":
                    score *= 2.5
                    score += 3.0

                # Boost image captions
                if any(w in query_lower for w in ["diagram", "figure", "chart", "graph", "plot", "image"]) and chunk.source_type == "image_caption":
                    score *= 2.5
                    score += 3.0

                if score > 0:
                    scored_chunks.append((score, chunk))

            scored_chunks.sort(key=lambda x: x[0], reverse=True)
            for _, c in scored_chunks[:top_k]:
                if not any(sc.chunk_id == c.chunk_id for sc in selected_chunks):
                    selected_chunks.append(c)

        # 3. Persistent SQLite FTS5 BM25 search fallback
        if len(selected_chunks) < top_k:
            fts_matches = store.search(doc_id=doc_id, query=query, limit=top_k)
            for m in fts_matches:
                if not any(sc.chunk_id == m["chunk_id"] for sc in selected_chunks):
                    selected_chunks.append(
                        DocumentChunk(
                            chunk_id=m["chunk_id"],
                            doc_id=m["doc_id"],
                            page=m["page"],
                            source_type=m["source_type"],
                            content=m["content"],
                        )
                    )

        if not selected_chunks and doc and doc.chunks:
            selected_chunks = doc.chunks[:3]

        context_blocks = []
        for c in selected_chunks[:top_k]:
            context_blocks.append(f"[Page {c.page} | Type: {c.source_type}]\n{c.content.strip()}")

        formatted_context = "\n\n".join(context_blocks)


        # Check doc status
        status_note = ""
        if doc and doc.status in {"processing_text", "text_ready", "processing_enrichment"}:
            status_note = (
                "Note: This document is still processing its tables/images — "
                "I can answer from the text for now, and give you a fuller answer in a moment."
            )

        metadata = {
            "doc_id": doc_id,
            "status": doc.status if doc else "unknown",
            "total_chunks": len(doc.chunks) if doc else 0,
            "retrieved_count": len(selected_chunks),
            "source_types_retrieved": list({c.source_type for c in selected_chunks}),
        }

        return formatted_context, status_note, metadata


# Global singleton instance
doc_processor = DocumentProcessor()
